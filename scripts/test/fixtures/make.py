#!/usr/bin/env python3
"""Regenerate the extraction test fixtures. See README.md in this directory."""
from pathlib import Path
from pptx import Presentation
from docx import Document

HERE = Path(__file__).parent

prs = Presentation()
for title, bullets, notes in [
    ("Vikat vs. Legacy CSPM",
     ["Where we win", "Agent-aware policy", "MCP server discovery"],
     "Lead with agent visibility. Do not claim they have no MCP coverage - they announced beta support."),
    ("Pricing bands",
     ["List: per-node", "Floor requires Deal Desk", "3-year commit unlocks band C"],
     "Never quote floor without approval."),
    ("Reference customers", ["Two public logos", "Four under NDA"], ""),
]:
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    tf = s.placeholders[1].text_frame
    tf.text = bullets[0]
    for b in bullets[1:]:
        tf.add_paragraph().text = b
    if notes:
        s.notes_slide.notes_text_frame.text = notes
prs.save(HERE / "battlecard.pptx")

doc = Document()
doc.add_heading("Deployment Models", level=1)
doc.add_paragraph("Vikat deploys as SaaS, in-VPC, or fully self-hosted for air-gapped environments.")
doc.add_paragraph("Self-hosted requires Kubernetes 1.27 or later.")
doc.add_heading("Engagement Models", level=1)
doc.add_paragraph("Standard pilot is six weeks, scoped to one solution area.")
doc.add_heading("Ampersands & entities", level=2)
doc.add_paragraph("Text with <angle brackets> and & ampersands should survive.")
doc.save(HERE / "deployment.docx")

print("fixtures regenerated")
